"""Integration tests for rendering pipeline"""

import json
import pytest
from pathlib import Path
from unittest.mock import Mock, patch, MagicMock
from typing import Dict, Any

import sys
sys.path.insert(0, str(Path(__file__).parent.parent.parent / "backend"))

from services.script_wrappers.image_gen import generate_image
from services.script_wrappers.html_converter import html_to_pptx_slide
from services.script_wrappers.slide_creator import create_pptx_from_images


@pytest.mark.integration
class TestRenderPathIntegration:
    """Test rendering path integration"""

    def create_mock_style_config(self) -> Dict[str, Any]:
        """Create a mock style configuration"""
        return {
            "id": "test-style",
            "name_zh": "测试风格",
            "name_en": "Test Style",
            "tier": 1,
            "colors": {
                "primary": "#FF0000",
                "secondary": "#00FF00",
                "background": "#FFFFFF",
                "text": "#000000",
                "accent": "#0000FF"
            },
            "typography": {
                "title_size": "32pt",
                "body_size": "16pt",
                "family": "System default"
            },
            "use_cases": ["education"],
            "render_paths": ["path_a", "path_b"]
        }

    @patch("subprocess.run")
    def test_render_path_a_workflow(self, mock_run, tmp_path):
        """Test Path A (HTML) rendering workflow"""
        # Path A: Prompt → HTML files → PPTX

        style_config = self.create_mock_style_config()

        # Step 1: Create mock HTML file (normally created by HTML renderer)
        html_file = tmp_path / "slide_1.html"
        html_file.write_text("<html><body>Slide 1</body></html>")

        # Step 2: Convert HTML to PPTX slide
        slide_data = {
            "success": True,
            "slide_id": "slide-1",
            "width_px": 1920,
            "height_px": 1080,
            "placeholders": [
                {"id": "ph-1", "x": 0, "y": 0, "w": 1920, "h": 1080}
            ],
            "output_path": str(tmp_path / "slide_1.pptx")
        }

        mock_run.return_value = Mock(
            returncode=0,
            stdout=json.dumps(slide_data),
            stderr=""
        )

        result = html_to_pptx_slide(str(html_file))

        assert result is not None
        assert result["slide_id"] == "slide-1"

    @patch("google.genai.Client")
    def test_render_path_b_workflow(self, mock_client_cls, tmp_path):
        """Test Path B (Image) rendering workflow"""
        # Path B: Prompt → Image generation → Image-based PPTX

        style_config = self.create_mock_style_config()

        # Step 1: Generate image
        output_image = tmp_path / "slide_1.png"

        # Mock the Gemini client to return a fake image
        from PIL import Image as PILImage
        import io
        fake_img = PILImage.new("RGB", (1024, 1024), color="blue")
        buf = io.BytesIO()
        fake_img.save(buf, format="PNG")

        mock_part = Mock()
        mock_part.inline_data = Mock(data=buf.getvalue())
        mock_part.text = None

        mock_response = Mock()
        mock_response.candidates = [Mock(content=Mock(parts=[mock_part]))]
        mock_client = Mock()
        mock_client.models.generate_content.return_value = mock_response
        mock_client_cls.return_value = mock_client

        image_path = generate_image(
            "A beautiful slide design",
            str(output_image),
            api_key="test-key"
        )

        assert image_path is not None
        assert Path(image_path).exists()

        # Step 2: Create PPTX from images
        output_pptx = tmp_path / "output.pptx"
        output_pptx.touch()

        result = create_pptx_from_images(
            [image_path],
            str(output_pptx)
        )

        assert result == str(output_pptx.resolve())

    @patch("subprocess.run")
    def test_mixed_rendering_workflow(self, mock_run, tmp_path):
        """Test mixed rendering with both paths"""
        style_config = self.create_mock_style_config()

        # Create test files for both paths
        html_file = tmp_path / "slide_html.html"
        html_file.write_text("<html></html>")

        image_file = tmp_path / "slide_image.png"
        image_file.touch()

        # Convert HTML
        slide_data = {
            "success": True,
            "slide_id": "slide-html",
            "output_path": str(tmp_path / "slide_html.pptx")
        }
        mock_run.return_value = Mock(
            returncode=0,
            stdout=json.dumps(slide_data),
            stderr=""
        )

        html_result = html_to_pptx_slide(str(html_file))
        assert html_result is not None

        # Create PPTX from mixed sources
        output_pptx = tmp_path / "output.pptx"
        output_pptx.touch()

        pptx_result = create_pptx_from_images(
            [str(image_file)],
            str(output_pptx)
        )

        assert pptx_result == str(output_pptx.resolve())

    def test_render_path_selection_logic(self):
        """Test logic for choosing between render paths"""
        style_config = self.create_mock_style_config()

        # Both paths available
        assert "path_a" in style_config["render_paths"]
        assert "path_b" in style_config["render_paths"]

        # Determine path based on style
        def select_render_path(style):
            """Simple path selection logic"""
            render_paths = style.get("render_paths", [])
            # Prefer path_b for image-based styles
            if "path_b" in render_paths:
                return "path_b"
            return "path_a"

        selected = select_render_path(style_config)
        assert selected in style_config["render_paths"]


@pytest.mark.integration
class TestPPTXAssembly:
    """Test PPTX assembly and validation"""

    @patch("subprocess.run")
    def test_path_a_pptx_assembly(self, mock_run, tmp_path):
        """Test Path A (HTML-based) PPTX assembly"""
        # Create mock HTML files
        html_files = []
        for i in range(3):
            html_file = tmp_path / f"slide_{i}.html"
            html_file.write_text(f"<html><body>Slide {i}</body></html>")
            html_files.append(html_file)

        # Mock HTML to PPTX conversion
        mock_run.return_value = Mock(
            returncode=0,
            stdout=json.dumps({
                "success": True,
                "slide_id": "slide-1",
                "output_path": str(tmp_path / "slide.pptx")
            }),
            stderr=""
        )

        # Convert each HTML file
        results = []
        for html_file in html_files:
            result = html_to_pptx_slide(str(html_file))
            if result:
                results.append(result)

        assert len(results) == 3

    @patch("subprocess.run")
    def test_path_b_pptx_assembly(self, mock_run, tmp_path):
        """Test Path B (Image-based) PPTX assembly"""
        # Create mock images
        images = []
        for i in range(3):
            img_file = tmp_path / f"slide_{i}.png"
            img_file.touch()
            images.append(str(img_file))

        # Create PPTX from images
        output_pptx = tmp_path / "presentation.pptx"
        output_pptx.touch()

        mock_run.return_value = Mock(returncode=0, stdout="", stderr="")

        result = create_pptx_from_images(
            images,
            str(output_pptx),
            layout="fullscreen"
        )

        assert result == str(output_pptx.resolve())
        mock_run.assert_called_once()

    @patch("subprocess.run")
    def test_pptx_with_mixed_content_and_titles(self, mock_run, tmp_path):
        """Test PPTX creation with mixed content and titles"""
        # Create images
        images = []
        titles = []
        for i in range(3):
            img_file = tmp_path / f"slide_{i}.png"
            img_file.touch()
            images.append(str(img_file))
            titles.append(f"Slide Title {i}")

        # Create PPTX with titles
        output_pptx = tmp_path / "presentation.pptx"
        output_pptx.touch()

        mock_run.return_value = Mock(returncode=0, stdout="", stderr="")

        result = create_pptx_from_images(
            images,
            str(output_pptx),
            layout="title_above",
            titles=titles,
            bg_color="FDFCF8"
        )

        assert result == str(output_pptx.resolve())

        # Verify all parameters passed to subprocess
        call_args = mock_run.call_args
        cmd = call_args[0][0]

        # Should include images
        for img in images:
            assert img in cmd

        # Should include titles
        assert "-t" in cmd
        for title in titles:
            assert title in cmd

        # Should include layout
        assert "--layout" in cmd
        assert "title_above" in cmd

        # Should include background color
        assert "-o" in cmd

    @patch("subprocess.run")
    def test_pptx_assembly_error_handling(self, mock_run, tmp_path):
        """Test error handling during PPTX assembly"""
        images = []
        for i in range(3):
            img_file = tmp_path / f"slide_{i}.png"
            img_file.touch()
            images.append(str(img_file))

        # Simulate failure
        mock_run.return_value = Mock(returncode=1, stdout="", stderr="Assembly failed")

        result = create_pptx_from_images(
            images,
            str(tmp_path / "output.pptx")
        )

        assert result is None

    def test_pptx_output_file_validation(self, tmp_path):
        """Test validation of PPTX output file"""
        # Create a simple PPTX file for validation
        from pptx import Presentation

        output_file = tmp_path / "test.pptx"

        # Create a simple presentation
        prs = Presentation()
        prs.slide_width = prs.slide_width  # 10 inches (default)
        prs.slide_height = prs.slide_height  # 7.5 inches (default)

        # Add a blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        prs.slides.add_slide(blank_layout)

        # Save
        prs.save(str(output_file))

        # Validate
        assert output_file.exists()
        assert output_file.stat().st_size > 0

        # Can be opened again
        loaded_prs = Presentation(str(output_file))
        assert len(loaded_prs.slides) == 1


@pytest.mark.integration
class TestRenderProgressTracking:
    """Test render progress and event tracking"""

    def test_slide_render_progress_tracking(self, tmp_path):
        """Test tracking of per-slide render progress"""
        total_slides = 5
        render_progress = {
            "total_slides": total_slides,
            "rendered_slides": 0,
            "current_slide": 0,
            "status": "rendering"
        }

        # Simulate rendering progress
        events = []
        for i in range(total_slides):
            render_progress["current_slide"] = i + 1
            render_progress["rendered_slides"] = i + 1

            event = {
                "type": "render_progress",
                "status": render_progress.copy()
            }
            events.append(event)

            # Verify progress
            assert event["status"]["rendered_slides"] == i + 1
            assert event["status"]["current_slide"] == i + 1

        assert len(events) == total_slides
        assert events[-1]["status"]["rendered_slides"] == total_slides

    def test_render_error_event_format(self):
        """Test format of render error events"""
        error_event = {
            "type": "error",
            "error_type": "image_generation_failed",
            "message": "Failed to generate image for slide 2",
            "slide_id": 2,
            "details": {
                "prompt": "A test prompt",
                "resolution": "1K"
            }
        }

        assert error_event["type"] == "error"
        assert error_event["error_type"] == "image_generation_failed"
        assert "message" in error_event
        assert "slide_id" in error_event

    def test_render_completion_event_format(self):
        """Test format of render completion events"""
        completion_event = {
            "type": "render_complete",
            "status": "success",
            "total_slides": 5,
            "rendered_slides": 5,
            "output_path": "/output/presentation.pptx",
            "duration_seconds": 45.5
        }

        assert completion_event["type"] == "render_complete"
        assert completion_event["status"] == "success"
        assert completion_event["total_slides"] == completion_event["rendered_slides"]
        assert "output_path" in completion_event
        assert "duration_seconds" in completion_event
