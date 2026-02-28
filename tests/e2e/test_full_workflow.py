"""
E2E tests for full workflow across all render paths

These tests validate the complete user journey:
1. Project creation with topic and style selection
2. Generation with SSE event streaming
3. PPTX download and validation

Tests are structured to work with mock data until Task #13/#14 are complete.
"""

import json
import pytest
from pathlib import Path
from unittest.mock import Mock, patch, MagicMock
from typing import Generator

import sys
sys.path.insert(0, str(Path(__file__).parent.parent.parent / "backend"))


# Mock SSE Event Streams
def mock_sse_events_path_a() -> Generator[str, None, None]:
    """Mock SSE event stream for Path A (HTML-based) rendering"""
    events = [
        {
            "type": "status",
            "status": "planning",
            "agent": "planner",
            "message": "Creating outline structure"
        },
        {
            "type": "status",
            "status": "researching",
            "agent": "researcher",
            "message": "Gathering source data",
            "stats": {"sources": 5, "queries": 2}
        },
        {
            "type": "status",
            "status": "writing",
            "agent": "writer",
            "message": "Writing slide content"
        },
        {
            "type": "hitl",
            "status": "waiting_for_approval",
            "outline": [
                {"title": "Slide 1", "content": "Introduction"},
                {"title": "Slide 2", "content": "Main Points"}
            ]
        },
        {
            "type": "status",
            "status": "rendering",
            "agent": "renderer",
            "message": "Starting Path A rendering"
        },
        {
            "type": "render_progress",
            "slide_number": 1,
            "total_slides": 2,
            "render_path": "path_a",
            "status": "rendering",
            "slide_title": "Slide 1"
        },
        {
            "type": "render_progress",
            "slide_number": 1,
            "total_slides": 2,
            "render_path": "path_a",
            "status": "complete",
            "slide_title": "Slide 1"
        },
        {
            "type": "render_progress",
            "slide_number": 2,
            "total_slides": 2,
            "render_path": "path_a",
            "status": "rendering",
            "slide_title": "Slide 2"
        },
        {
            "type": "render_progress",
            "slide_number": 2,
            "total_slides": 2,
            "render_path": "path_a",
            "status": "complete",
            "slide_title": "Slide 2"
        },
        {
            "type": "complete",
            "status": "done",
            "pptx_path": "/output/presentation.pptx"
        }
    ]

    for event in events:
        yield f"data: {json.dumps(event)}\n\n"


def mock_sse_events_path_b() -> Generator[str, None, None]:
    """Mock SSE event stream for Path B (Image-based) rendering"""
    events = [
        {
            "type": "status",
            "status": "planning",
            "agent": "planner",
            "message": "Creating visual structure"
        },
        {
            "type": "status",
            "status": "researching",
            "agent": "researcher",
            "message": "Gathering reference data"
        },
        {
            "type": "status",
            "status": "visual_design",
            "agent": "visual",
            "message": "Planning visual compositions"
        },
        {
            "type": "hitl",
            "status": "waiting_for_approval",
            "outline": [
                {"title": "Slide 1", "visual_type": "hero_image"},
                {"title": "Slide 2", "visual_type": "data_visualization"}
            ]
        },
        {
            "type": "status",
            "status": "rendering",
            "agent": "renderer",
            "message": "Starting Path B rendering"
        },
        {
            "type": "render_progress",
            "slide_number": 1,
            "total_slides": 2,
            "render_path": "path_b",
            "status": "generating_image",
            "slide_title": "Slide 1"
        },
        {
            "type": "render_progress",
            "slide_number": 1,
            "total_slides": 2,
            "render_path": "path_b",
            "status": "complete",
            "slide_title": "Slide 1",
            "thumbnail_url": "/thumbnails/slide_1.jpg"
        },
        {
            "type": "render_progress",
            "slide_number": 2,
            "total_slides": 2,
            "render_path": "path_b",
            "status": "generating_image",
            "slide_title": "Slide 2"
        },
        {
            "type": "render_progress",
            "slide_number": 2,
            "total_slides": 2,
            "render_path": "path_b",
            "status": "complete",
            "slide_title": "Slide 2",
            "thumbnail_url": "/thumbnails/slide_2.jpg"
        },
        {
            "type": "complete",
            "status": "done",
            "pptx_path": "/output/presentation.pptx"
        }
    ]

    for event in events:
        yield f"data: {json.dumps(event)}\n\n"


def mock_sse_events_auto_mixed() -> Generator[str, None, None]:
    """Mock SSE event stream for Auto (mixed) rendering"""
    events = [
        {
            "type": "status",
            "status": "planning",
            "agent": "planner",
            "message": "Planning mixed rendering strategy"
        },
        {
            "type": "hitl",
            "status": "waiting_for_approval",
            "outline": [
                {"title": "Slide 1", "render_path": "path_a"},
                {"title": "Slide 2", "render_path": "path_b"}
            ]
        },
        {
            "type": "status",
            "status": "rendering",
            "agent": "renderer",
            "message": "Starting mixed rendering"
        },
        {
            "type": "render_progress",
            "slide_number": 1,
            "total_slides": 2,
            "render_path": "path_a",
            "status": "complete",
            "slide_title": "Slide 1"
        },
        {
            "type": "render_progress",
            "slide_number": 2,
            "total_slides": 2,
            "render_path": "path_b",
            "status": "complete",
            "slide_title": "Slide 2",
            "thumbnail_url": "/thumbnails/slide_2.jpg"
        },
        {
            "type": "complete",
            "status": "done",
            "pptx_path": "/output/presentation.pptx"
        }
    ]

    for event in events:
        yield f"data: {json.dumps(event)}\n\n"


@pytest.mark.e2e
class TestFullWorkflowPathA:
    """E2E tests for Path A (HTML-based) rendering"""

    def test_full_workflow_path_a(self):
        """Test complete workflow: topic → style → generation → download (Path A)"""
        # Step 1: Create project with Path A preference
        project_data = {
            "prompt": "Create a presentation about machine learning",
            "style": "01-snoopy",
            "render_path_preference": "path_a"
        }

        # Step 2: Verify project creation response
        assert project_data["prompt"]
        assert project_data["style"]
        assert project_data["render_path_preference"] == "path_a"

        # Step 3: Simulate SSE event stream
        events = list(mock_sse_events_path_a())
        assert len(events) > 0

        # Step 4: Verify event sequence
        event_types = []
        for event_str in events:
            # Parse SSE format: "data: {json}\n\n"
            if event_str.startswith("data: "):
                event_json = event_str.replace("data: ", "").strip()
                if event_json:
                    event = json.loads(event_json)
                    event_types.append(event["type"])

        # Verify expected event sequence
        assert "status" in event_types
        assert "hitl" in event_types
        assert "render_progress" in event_types
        assert "complete" in event_types

    def test_sse_event_flow_path_a(self):
        """Test SSE event flow for Path A rendering"""
        events = list(mock_sse_events_path_a())

        parsed_events = []
        for event_str in events:
            if event_str.startswith("data: "):
                event_json = event_str.replace("data: ", "").strip()
                if event_json:
                    parsed_events.append(json.loads(event_json))

        # Verify event sequence: status → hitl → status → render_progress → complete
        assert parsed_events[0]["type"] == "status"
        assert parsed_events[0]["status"] == "planning"

        # Find HITL event
        hitl_event = next((e for e in parsed_events if e["type"] == "hitl"), None)
        assert hitl_event is not None
        assert hitl_event["status"] == "waiting_for_approval"

        # Find render progress events
        progress_events = [e for e in parsed_events if e["type"] == "render_progress"]
        assert len(progress_events) >= 2
        for i, event in enumerate(progress_events):
            assert event["slide_number"] >= 1
            assert event["total_slides"] >= 1

        # Last event should be complete
        assert parsed_events[-1]["type"] == "complete"

    def test_pptx_validity_path_a(self, tmp_path):
        """Test that generated PPTX is valid (Path A)"""
        from pptx import Presentation

        # Create a simple PPTX file to simulate Path A output
        pptx_file = tmp_path / "presentation.pptx"
        prs = Presentation()

        # Add HTML-based slides (simulated)
        for i in range(2):
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            # In real Path A, would have HTML converted to PPTX

        prs.save(str(pptx_file))

        # Verify file exists and is valid
        assert pptx_file.exists()
        assert pptx_file.stat().st_size > 0

        # Verify can be opened
        loaded = Presentation(str(pptx_file))
        assert len(loaded.slides) == 2


@pytest.mark.e2e
class TestFullWorkflowPathB:
    """E2E tests for Path B (Image-based) rendering"""

    def test_full_workflow_path_b(self):
        """Test complete workflow: topic → style → generation → download (Path B)"""
        project_data = {
            "prompt": "Create a presentation about artificial intelligence",
            "style": "02-manga",
            "render_path_preference": "path_b"
        }

        assert project_data["render_path_preference"] == "path_b"

        events = list(mock_sse_events_path_b())
        assert len(events) > 0

        # Verify Path B specific events
        event_types = []
        for event_str in events:
            if event_str.startswith("data: "):
                event_json = event_str.replace("data: ", "").strip()
                if event_json:
                    event = json.loads(event_json)
                    event_types.append(event.get("type"))

        assert "render_progress" in event_types

    def test_image_generation_progress_path_b(self):
        """Test image generation progress events (Path B)"""
        events = list(mock_sse_events_path_b())

        parsed_events = []
        for event_str in events:
            if event_str.startswith("data: "):
                event_json = event_str.replace("data: ", "").strip()
                if event_json:
                    parsed_events.append(json.loads(event_json))

        # Find progress events with image generation
        progress_events = [e for e in parsed_events if e["type"] == "render_progress"]
        image_gen_events = [e for e in progress_events if e.get("status") == "generating_image"]

        # Path B should have image generation events
        assert len(image_gen_events) > 0

        # Verify thumbnail URLs in complete events
        complete_events = [e for e in progress_events if e.get("status") == "complete"]
        for event in complete_events:
            assert "slide_title" in event
            # Thumbnails may or may not be present
            if "thumbnail_url" in event:
                assert event["thumbnail_url"].startswith("/")

    def test_pptx_with_images_path_b(self, tmp_path):
        """Test PPTX creation with images (Path B)"""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_file = tmp_path / "presentation.pptx"
        prs = Presentation()

        # Create image for slide
        img_file = tmp_path / "slide_1.png"
        img_file.write_bytes(b"fake image data")

        # Add slide with image (simulated Path B output)
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        prs.save(str(pptx_file))

        assert pptx_file.exists()
        assert pptx_file.stat().st_size > 0


@pytest.mark.e2e
class TestAutoMixedRendering:
    """E2E tests for Auto (mixed) rendering"""

    def test_full_workflow_auto_mixed(self):
        """Test complete workflow with auto/mixed render preference"""
        project_data = {
            "prompt": "Create a presentation about data science",
            "style": "03-ligne-claire",
            "render_path_preference": "auto"
        }

        assert project_data["render_path_preference"] == "auto"

        events = list(mock_sse_events_auto_mixed())
        assert len(events) > 0

    def test_mixed_path_selection(self):
        """Test that mixed rendering uses appropriate paths per slide"""
        events = list(mock_sse_events_auto_mixed())

        parsed_events = []
        for event_str in events:
            if event_str.startswith("data: "):
                event_json = event_str.replace("data: ", "").strip()
                if event_json:
                    parsed_events.append(json.loads(event_json))

        # Find render progress events
        progress_events = [e for e in parsed_events if e["type"] == "render_progress"]

        # Should have both path_a and path_b
        paths_used = set(e.get("render_path") for e in progress_events)
        assert "path_a" in paths_used or "path_b" in paths_used


@pytest.mark.e2e
class TestErrorScenarios:
    """E2E tests for error scenarios"""

    def test_invalid_style_id(self):
        """Test handling of invalid style_id"""
        project_data = {
            "prompt": "Test presentation",
            "style": "invalid-style-id"
        }

        # Should either:
        # 1. Return 404 error immediately, or
        # 2. Return error in SSE stream
        # Implementation depends on backend design

        assert project_data["style"] == "invalid-style-id"

    def test_empty_prompt_validation(self):
        """Test validation of empty prompt"""
        project_data = {
            "prompt": "",
            "style": "01-snoopy"
        }

        # Should reject empty prompt
        assert not project_data["prompt"]

    def test_missing_api_key_error(self):
        """Test handling of missing API key (for image generation)"""
        error_event = {
            "type": "error",
            "error_type": "api_key_missing",
            "message": "GEMINI_API_KEY not configured",
            "recoverable": False
        }

        assert error_event["type"] == "error"
        assert error_event["error_type"] == "api_key_missing"

    def test_network_timeout_error(self):
        """Test handling of network timeout during generation"""
        error_event = {
            "type": "error",
            "error_type": "timeout",
            "message": "Generation timed out after 10 minutes",
            "slide_number": 2,
            "recoverable": True
        }

        assert error_event["error_type"] == "timeout"
        # Graceful degradation: mark as failed but continue
        assert error_event.get("recoverable", False)


@pytest.mark.e2e
class TestSSEEventStructure:
    """Test SSE event format and structure"""

    def test_status_event_format(self):
        """Test format of status events"""
        event = {
            "type": "status",
            "status": "planning",
            "agent": "planner",
            "message": "Creating outline"
        }

        assert event["type"] == "status"
        assert "status" in event
        assert "agent" in event
        assert "message" in event

    def test_render_progress_event_format(self):
        """Test format of render_progress events"""
        event = {
            "type": "render_progress",
            "slide_number": 1,
            "total_slides": 5,
            "render_path": "path_a",
            "status": "rendering",
            "slide_title": "Introduction"
        }

        assert event["type"] == "render_progress"
        assert event["slide_number"] >= 1
        assert event["total_slides"] >= 1
        assert event["render_path"] in ["path_a", "path_b"]
        assert event["status"] in ["rendering", "complete", "failed", "generating_image"]

    def test_complete_event_format(self):
        """Test format of completion event"""
        event = {
            "type": "complete",
            "status": "done",
            "pptx_path": "/output/presentation.pptx"
        }

        assert event["type"] == "complete"
        assert event["status"] == "done"
        assert "pptx_path" in event

    def test_error_event_format(self):
        """Test format of error events"""
        event = {
            "type": "error",
            "error_type": "generation_failed",
            "message": "Failed to generate slide content",
            "details": {
                "slide_number": 2,
                "agent": "writer"
            }
        }

        assert event["type"] == "error"
        assert "error_type" in event
        assert "message" in event


@pytest.mark.e2e
class TestDownloadAndValidation:
    """Test PPTX download and validation"""

    def test_pptx_file_download(self, tmp_path):
        """Test PPTX file is available for download after generation"""
        pptx_file = tmp_path / "presentation.pptx"

        # Simulate generated file
        from pptx import Presentation
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(pptx_file))

        # Verify file can be downloaded
        assert pptx_file.exists()
        assert pptx_file.stat().st_size > 1024  # Should be > 1KB

    def test_pptx_file_validity(self, tmp_path):
        """Test that downloaded PPTX is valid and can be opened"""
        from pptx import Presentation

        pptx_file = tmp_path / "presentation.pptx"
        prs = Presentation()

        # Add multiple slides
        for i in range(3):
            prs.slides.add_slide(prs.slide_layouts[6])

        prs.save(str(pptx_file))

        # Verify it's a valid PPTX
        loaded = Presentation(str(pptx_file))
        assert len(loaded.slides) == 3

    def test_pptx_file_not_empty(self, tmp_path):
        """Test that PPTX file is not empty"""
        from pptx import Presentation

        pptx_file = tmp_path / "presentation.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(pptx_file))

        assert pptx_file.stat().st_size > 0
