"""
渲染引擎 Agent (Renderer) 工具函数
包含 python-pptx 操作工具
"""

from pathlib import Path
from typing import Dict, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# 输出目录
OUTPUT_DIR = Path(__file__).parent.parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

# 布局映射
LAYOUT_MAP = {
    "title_slide": 0,
    "cover": 0,
    "bullet_list": 1,
    "bullet_points": 1,
    "content": 1,
    "two_content": 3,
    "two_column": 3,
    "comparison": 4,
    "picture_with_caption": 8,
    "blank": 6,
    "conclusion": 1,
    "data": 1,
    "chart": 1,
    "quote": 6,
    "data_driven": 1,
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    """将十六进制颜色转换为 RGB"""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def apply_theme_to_slide(slide, theme_config: Dict):
    """应用主题样式到幻灯片"""
    colors = theme_config.get("colors", {})
    bg_color = colors.get("background", "#FDFCF8")
    
    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(bg_color)


def add_text_to_placeholder(placeholder, text: str, theme_config: Dict, is_title: bool = False):
    """向占位符添加文本"""
    colors = theme_config.get("colors", {})
    text_color = colors.get("text", "#2C2C24")
    primary_color = colors.get("primary", "#5D7052")
    
    tf = placeholder.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(44 if is_title else 24)
    p.font.color.rgb = hex_to_rgb(primary_color if is_title else text_color)
    p.font.bold = is_title


def add_bullet_points(placeholder, points: list, theme_config: Dict):
    """添加要点列表"""
    colors = theme_config.get("colors", {})
    text_color = colors.get("text", "#2C2C24")
    
    tf = placeholder.text_frame
    tf.clear()
    
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = str(point)
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(text_color)
        p.level = 0


def get_layout_id(layout_name: str, total_layouts: int) -> int:
    """获取布局 ID，确保在有效范围内"""
    layout_id = LAYOUT_MAP.get(layout_name, 1)
    
    # 确保 layout_id 在有效范围内
    if layout_id >= total_layouts:
        layout_id = min(1, total_layouts - 1)
    
    return layout_id


def create_presentation() -> Presentation:
    """创建新的演示文稿"""
    prs = Presentation()
    
    # 设置幻灯片尺寸 (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    return prs


def save_presentation(prs: Presentation, session_id: str) -> str:
    """保存演示文稿"""
    filename = f"presentation_{session_id}.pptx"
    filepath = OUTPUT_DIR / filename
    prs.save(str(filepath))
    return str(filepath)
