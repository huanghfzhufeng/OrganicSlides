"""
PPTX Assembler — combines per-slide render results into a single presentation.

Handles three cases:
  - All Path A results: merge individual single-slide PPTX files
  - All Path B results: assemble image files via python-pptx directly
  - Mixed: merge both types into one presentation

The assembler never mutates the input list; it produces a new PPTX file.
"""

from __future__ import annotations

import logging
import shutil
import uuid
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.util import Inches

if TYPE_CHECKING:
    from agents.renderer.paths import SlideRenderResult

logger = logging.getLogger(__name__)

# 16:9 slide dimensions (standard)
_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def assemble_presentation(
    slide_results: list["SlideRenderResult"],
    output_path: str,
) -> str:
    """
    Merge all slide render results into a single PPTX file at *output_path*.

    - Path A results (*.pptx files) are merged slide-by-slide via python-pptx.
    - Path B results (*.png files) are placed fullscreen on blank slides.
    - Failed slides are skipped with a warning; a placeholder slide is inserted.

    Returns the absolute path to the assembled PPTX on success.
    Raises RuntimeError if no successful slides exist.
    """
    successful = [r for r in slide_results if r.success]
    if not successful:
        raise RuntimeError("No slides rendered successfully — cannot assemble presentation")

    prs = _new_presentation()

    for result in sorted(successful, key=lambda r: r.slide_index):
        try:
            # path_a, path_a_fallback, fallback_text → PPTX file merge
            # path_b → PNG image slide
            if result.render_path in ("path_a", "path_a_fallback", "fallback_text"):
                _merge_pptx_slide(prs, result.output_path)
            else:
                _add_image_slide(prs, result.output_path)
        except Exception as exc:
            logger.error(
                "Failed to add slide %d (%s): %s",
                result.slide_index, result.render_path, exc,
            )
            _add_error_placeholder_slide(prs, result.slide_index, str(exc))

    # Log any failed slides
    failed = [r for r in slide_results if not r.success]
    for r in failed:
        logger.warning(
            "Slide %d was not rendered (error: %s)",
            r.slide_index, r.error,
        )
        _add_error_placeholder_slide(prs, r.slide_index, r.error or "render failed")

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    logger.info("Assembled %d slides → %s", len(prs.slides), out)
    return str(out.resolve())


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = _SLIDE_WIDTH
    prs.slide_height = _SLIDE_HEIGHT
    return prs


def _merge_pptx_slide(target_prs: Presentation, source_pptx_path: str) -> None:
    """
    Copy the first slide from *source_pptx_path* into *target_prs*.

    python-pptx does not support direct slide merging, so we use the
    xml-level approach: clone the slide XML and copy all referenced blobs.
    """
    src = Presentation(source_pptx_path)
    if not src.slides:
        logger.warning("Source PPTX has no slides: %s", source_pptx_path)
        return

    # Use blank layout as carrier — then replace its XML
    blank_layout = _get_blank_layout(target_prs)
    target_slide = target_prs.slides.add_slide(blank_layout)

    src_slide = src.slides[0]

    # Copy background
    _copy_slide_background(src_slide, target_slide, src, target_prs)

    # Copy all shapes
    for shape in src_slide.shapes:
        _clone_shape(shape, target_slide, src, target_prs)


def _add_image_slide(prs: Presentation, image_path: str) -> None:
    """Add a fullscreen image slide to *prs*."""
    blank_layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(
        image_path,
        left=0,
        top=0,
        width=_SLIDE_WIDTH,
        height=_SLIDE_HEIGHT,
    )


def _add_error_placeholder_slide(prs: Presentation, slide_index: int, error: str) -> None:
    """Insert a clearly-labelled error slide so the deck is still complete."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    blank_layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(blank_layout)

    # Red background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEB)

    # Error text box
    txb = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.333), Inches(2.5)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"⚠ Slide {slide_index + 1} render error:\n{error[:300]}"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)


def _get_blank_layout(prs: Presentation):
    """Return the blank slide layout (index 6 by convention)."""
    idx = min(6, len(prs.slide_layouts) - 1)
    return prs.slide_layouts[idx]


# ---------------------------------------------------------------------------
# XML-level slide copy helpers (minimal implementation)
# ---------------------------------------------------------------------------

def _copy_slide_background(src_slide, target_slide, src_prs, target_prs) -> None:
    """Copy slide background fill from source to target."""
    try:
        src_bg = src_slide.background
        tgt_bg = target_slide.background

        if src_bg.fill.type is not None:
            # Replicate solid color fills
            from pptx.enum.dml import MSO_THEME_COLOR
            fill_type = src_bg.fill.type
            # Access background fill type integer directly
            if str(fill_type) == "SOLID (1)":
                tgt_bg.fill.solid()
                tgt_bg.fill.fore_color.rgb = src_bg.fill.fore_color.rgb
    except Exception as exc:
        logger.debug("Could not copy background: %s", exc)


def _clone_shape(src_shape, target_slide, src_prs, target_prs) -> None:
    """
    Clone a shape from src_slide into target_slide using XML copy.

    This is the standard low-level technique when python-pptx doesn't
    expose a higher-level copy API.
    """
    import copy
    from lxml import etree
    from pptx.oxml.ns import qn

    # Deep-copy the shape XML element
    sp_elem = copy.deepcopy(src_shape._element)

    # Handle embedded images: copy image blobs to target package
    if src_shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE == 13
        try:
            _copy_picture_blob(src_shape, sp_elem, src_prs, target_slide, target_prs)
        except Exception as exc:
            logger.debug("Could not copy picture blob: %s", exc)
            return

    # Append shape XML to target slide's spTree
    sp_tree = target_slide.shapes._spTree
    sp_tree.append(sp_elem)


def _copy_picture_blob(src_shape, sp_elem, src_prs, target_slide, target_prs) -> None:
    """Copy the image blob from source package to target package."""
    from lxml import etree
    from pptx.oxml.ns import qn

    # Find the rId for the image in src slide
    blipfill = sp_elem.find('.//' + qn('p:blipFill'))
    if blipfill is None:
        return
    blip = blipfill.find(qn('a:blip'))
    if blip is None:
        return

    r_embed = blip.get(qn('r:embed'))
    if not r_embed:
        return

    # Get image part from source slide
    src_slide = src_shape.part
    try:
        image_part = src_slide.related_parts[r_embed]
    except KeyError:
        return

    # Add image to target slide and get new rId
    from pptx.parts.image import ImagePart
    image_bytes = image_part.blob
    content_type = image_part.content_type

    new_image_part, new_rId = target_slide.part.get_or_add_image_part(
        image_bytes
    )
    # Update rId in the cloned element
    blip.set(qn('r:embed'), new_rId)
